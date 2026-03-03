"""Tools for finding and extracting figures from influential academic papers.

Paper search uses the Semantic Scholar REST API (free, no key required for
basic usage) with automatic fallback to OpenAlex when S2 is rate-limited.
Figure extraction uses PyMuPDF.  PPTX export uses python-pptx.

Output layout (all under workspace):
    outputs/figure_refs/<paper_id>/   PNG + SVG figures, paper.pdf
    outputs/figure_refs/reference_pack_<title>_<ts>.pptx
"""

from __future__ import annotations

import asyncio
import csv
import json
import re
import unicodedata
from contextlib import asynccontextmanager
from datetime import datetime
from pathlib import Path
from typing import Any

from loguru import logger

from clawphd.agent.tools.base import Tool


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

def _slugify(text: str, max_len: int = 60) -> str:
    """Convert arbitrary text to a filesystem-safe ASCII slug."""
    text = unicodedata.normalize("NFKD", text)
    text = re.sub(r"[^\w\s-]", "", text, flags=re.ASCII).strip().lower()
    text = re.sub(r"[\s_-]+", "_", text)
    return text[:max_len]


def _timestamp() -> str:
    return datetime.now().strftime("%Y%m%d_%H%M%S")


# ---------------------------------------------------------------------------
# Semantic Scholar helpers
# ---------------------------------------------------------------------------

_S2_SEARCH = "https://api.semanticscholar.org/graph/v1/paper/search"
_S2_REFS = "https://api.semanticscholar.org/graph/v1/paper/{pid}/references"
_S2_FIELDS = "paperId,externalIds,title,year,citationCount,openAccessPdf,isOpenAccess"


@asynccontextmanager
async def _s2_client(api_key: str | None = None):
    """Async HTTP client pre-configured for Semantic Scholar."""
    import httpx

    headers: dict[str, str] = {"Accept": "application/json"}
    if api_key:
        headers["x-api-key"] = api_key
    async with httpx.AsyncClient(headers=headers, timeout=30.0, follow_redirects=True) as client:
        yield client


async def _s2_get(
    client: Any,
    url: str,
    params: dict | None = None,
    max_retries: int = 5,
) -> Any:
    """GET a Semantic Scholar endpoint with exponential back-off on 429 / 5xx."""
    for attempt in range(max_retries):
        resp = await client.get(url, params=params)
        if resp.status_code < 400:
            return resp
        is_rate = resp.status_code == 429
        is_server = resp.status_code >= 500
        if not (is_rate or is_server) or attempt == max_retries - 1:
            resp.raise_for_status()
            return resp
        wait = int(resp.headers.get("Retry-After", 2 ** attempt))
        label = "rate limit (429)" if is_rate else f"server error ({resp.status_code})"
        logger.warning(
            "Semantic Scholar {} — retrying in {}s (attempt {}/{})",
            label, wait, attempt + 1, max_retries,
        )
        await asyncio.sleep(wait)
    resp.raise_for_status()
    return resp


# ---------------------------------------------------------------------------
# OpenAlex fallback helper
# ---------------------------------------------------------------------------

async def _openalex_search(
    topic: str,
    num_papers: int,
    min_citations: int,
    year_start: int,
) -> list[dict]:
    """Search OpenAlex (no API key needed, generous rate limits)."""
    import httpx

    url = "https://api.openalex.org/works"
    params = {
        "search": topic,
        "filter": f"cited_by_count:>{min_citations},publication_year:>{year_start - 1}",
        "per-page": min(num_papers * 3, 50),
        "select": "id,title,publication_year,cited_by_count,open_access,primary_location,ids",
    }
    async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as client:
        resp = await client.get(url, params=params)
        resp.raise_for_status()
    data = resp.json().get("results") or []

    results: list[dict] = []
    for item in data:
        title = item.get("title") or ""
        year = item.get("publication_year")
        cit = item.get("cited_by_count", 0)
        if cit < min_citations:
            continue

        # Prefer open-access landing page PDF
        pdf_url: str | None = None
        oa = item.get("open_access") or {}
        if oa.get("oa_url"):
            pdf_url = oa["oa_url"]
        if not pdf_url:
            loc = item.get("primary_location") or {}
            pdf_url = loc.get("pdf_url") or loc.get("landing_page_url")

        # Extract arxiv id
        ids = item.get("ids") or {}
        arxiv_id: str | None = None
        if ids.get("arxiv"):
            arxiv_id = ids["arxiv"].split("/")[-1]
        elif pdf_url and "arxiv.org" in (pdf_url or ""):
            m = re.search(r"(\d{4}\.\d{4,5})", pdf_url)
            if m:
                arxiv_id = m.group(1)

        paper_id = arxiv_id or _slugify(title)[:20]
        results.append({
            "paper_id": paper_id,
            "arxiv_id": arxiv_id,
            "title": title,
            "year": year,
            "citation_count": cit,
            "pdf_url": pdf_url,
            "source": "openalex",
        })
        if len(results) >= num_papers:
            break

    return results


# ---------------------------------------------------------------------------
# PyMuPDF figure extraction helpers
# ---------------------------------------------------------------------------

_FIG_RE = re.compile(r"(?:Figure|Fig\.?)\s+(\d+)\s*[:.]", re.IGNORECASE)

# Minimum fraction of column area that image blocks must cover to be trusted
# as the actual figure (vs. small legend symbols like coloured squares).
_MIN_IMG_FILL = 0.08

# Maximum distance (in points) to search above/below caption for figure content.
_MAX_SEARCH_PTS = 550


def _detect_columns(blocks: list, pw: float) -> tuple[int, float]:
    """Return (num_columns, midpoint_x) by analysing text block positions.

    Only considers blocks that look like body text: must be wider than 8 % of
    the page width (filters out figure-label fragments) and must contain at
    least 4 words (filters out short labels like 'Expert Gate').  If at least
    one such block exists on each side of the page centre, the layout is
    two-column.
    """
    mid = pw / 2
    left = right = 0
    for blk in blocks:
        if blk["type"] != 0:
            continue
        x0, _, x1, _ = blk["bbox"]
        blk_w = x1 - x0
        if blk_w > pw * 0.60:      # full-width element → skip
            continue
        if blk_w < pw * 0.08:      # very narrow → figure label or tick → skip
            continue
        text = "".join(
            s["text"]
            for line in blk.get("lines", [])
            for s in line.get("spans", [])
        )
        if len(text.split()) < 4:  # too few words → label, not body text → skip
            continue
        if (x0 + x1) / 2 < mid:
            left += 1
        else:
            right += 1
    if left >= 1 and right >= 1:
        return 2, mid
    return 1, mid


def _column_x_bounds(
    cx0: float,
    cx1: float,
    num_cols: int,
    mid_x: float,
    pw: float,
    blocks: list | None = None,
    cy0: float = 0.0,
    page: Any = None,
) -> tuple[float, float]:
    """Return the (x_left, x_right) crop bounds for the figure's caption.

    Detects spanning figures in two-column papers via three checks (in order):
    1. Caption width > 55 % of page.
    2. Raster image crossing mid_x (not shadowed by a previous-figure caption).
    3. Vector drawing crossing mid_x that is BOTH:
       a. Not shadowed by a previous-figure caption in the same x-region.
       b. Centered within 15 % of page width from mid_x (i.e. the element is
          symmetric around the column boundary, not just a one-column element
          with a small left/right overhang that slightly crosses mid_x).

    Rules (b) prevents false spanning detection when a figure in one column
    has drawing elements that marginally extend past mid_x, or when a
    full-width figure from a different position on the same page has crossing
    drawings above the current figure's caption.
    """
    if num_cols == 1:
        return 0.0, pw

    # Caption spanning > 55 % of page → full-width figure.
    if cx1 - cx0 > pw * 0.55:
        return 0.0, pw

    # Collect bounding boxes of other figure captions above cy0 on this page.
    # Used to avoid attributing their figure's drawings to the current figure.
    prev_caps: list[tuple[float, float, float]] = []  # (x0, x1, y_bottom)
    if blocks:
        for ob in blocks:
            if ob["type"] != 0:
                continue
            ox0, oy0, ox1, oy1 = ob["bbox"]
            if oy1 >= cy0 or cy0 - oy1 > _MAX_SEARCH_PTS:
                continue
            txt = "".join(s["text"] for l in ob.get("lines", []) for s in l.get("spans", []))
            if _FIG_RE.search(txt):
                prev_caps.append((ox0, ox1, oy1 + 2))

    def _is_shadowed(rx0: float, rx1: float, ry0: float) -> bool:
        """Return True if (rx0,rx1,ry0) is attributable to a previous figure."""
        for pcx0, pcx1, pcy_bot in prev_caps:
            if ry0 >= pcy_bot:
                continue   # drawing is below that caption — not shadowed by it
            x_overlap = min(rx1, pcx1) - max(rx0, pcx0)
            # 15 % overlap with the previous caption width is enough evidence
            # that this drawing belongs to the neighbouring figure's region.
            if x_overlap > (pcx1 - pcx0) * 0.15:
                return True
        return False

    # Raster image block crossing mid_x.
    if blocks:
        for ob in blocks:
            if ob["type"] != 1:
                continue
            ox0, oy0, ox1, oy1 = ob["bbox"]
            if oy1 > cy0 + 10 or cy0 - oy0 > _MAX_SEARCH_PTS:
                continue
            if ox0 < mid_x and ox1 > mid_x and not _is_shadowed(ox0, ox1, oy0):
                return 0.0, pw

    # Vector drawing crossing mid_x — apply two additional guards:
    #   (a) column-aware shadow check (drawing not from a neighbouring figure).
    #   (b) draw_cx within 15 % of page width from mid_x: a truly spanning
    #       element is centred near the column boundary; a one-column element
    #       that just overshoots mid_x has its centre far from mid_x.
    span_tolerance = pw * 0.15   # 92 pt for a 612-pt wide page
    if page is not None:
        try:
            for drw in page.get_drawings():
                rect = drw.get("rect")
                if rect is None:
                    continue
                if rect.y1 > cy0 + 10 or cy0 - rect.y0 > _MAX_SEARCH_PTS:
                    continue
                if rect.width < 15:     # skip thin rules / column dividers
                    continue
                if not (rect.x0 < mid_x and rect.x1 > mid_x):
                    continue
                draw_cx = (rect.x0 + rect.x1) / 2
                if abs(draw_cx - mid_x) >= span_tolerance:
                    continue   # centred too far in one column — not spanning
                # Left edge must reach clearly into the left column (not just
                # marginally cross mid_x from a right-column element).
                if rect.x0 >= mid_x - pw * 0.10:
                    continue   # element barely overshoots — not a spanning figure
                if _is_shadowed(rect.x0, rect.x1, rect.y0):
                    continue   # belongs to a different figure above
                return 0.0, pw
        except Exception:
            pass   # older PyMuPDF versions may not support get_drawings

    cap_cx = (cx0 + cx1) / 2
    if cap_cx < mid_x:
        return 0.0, mid_x - 4      # left column
    return mid_x + 4, pw           # right column


def _figure_top(
    cy0: float,
    cy1: float,
    col_x0: float,
    col_x1: float,
    blocks: list,
    page: Any = None,
    num_cols: int = 1,
    mid_x: float = 0.0,
    pw: float = 0.0,
) -> float:
    """Estimate the y-coordinate where the figure starts above a caption.

    Strategy (in priority order)
    ----------------------------
    0. Vector drawing boundaries: the topmost drawing above the caption gives
       a reliable estimate for pure-vector figures (charts, diagrams).
    1. Raster image blocks (type==1) in the same column that end above the
       caption and cover a substantial fraction of the searched area.
    2. The bottom edge of the nearest "body text" block in the column above
       the caption.  Body text (paragraphs, headings) immediately precedes the
       figure area.  Axis labels / legend text are excluded via word-count and
       width heuristics.
    3. The bottom of the nearest other figure caption above this one.
    4. Fixed 300 pt fallback, capped at the page top.
    """
    col_width = col_x1 - col_x0
    max_dist = min(_MAX_SEARCH_PTS, cy0)

    # --- Pass 0: vector drawing bounding box ---
    # The topmost drawing path above the caption is the most direct signal for
    # where a pure-vector figure (bar chart, architecture diagram, …) starts.
    #
    # Per-drawing column-aware shadow check: if a drawing is above the bottom
    # edge of another figure's caption AND it horizontally overlaps that caption
    # significantly, the drawing belongs to a different figure — skip it.  This
    # is more accurate than a single page-wide y cutoff because it handles pages
    # where two figures sit side-by-side in different columns at the same y level.
    # Infer the column x-range of each previous figure to get maximum shadow
    # coverage.  For a two-column layout the caption text block is usually
    # narrower than the figure it labels, so using the full column range (rather
    # than the caption bbox) prevents drawings from adjacent figures that extend
    # to the column edge from leaking through the shadow check.
    prev_caps_0: list[tuple[float, float, float]] = []  # (shade_x0, shade_x1, y_bottom)
    for ob in blocks:
        if ob["type"] != 0:
            continue
        ox0, oy0, ox1, oy1 = ob["bbox"]
        if oy1 >= cy0 or cy0 - oy1 > max_dist:
            continue
        txt = "".join(s["text"] for l in ob.get("lines", []) for s in l.get("spans", []))
        if _FIG_RE.search(txt):
            if pw > 0 and mid_x > 0 and num_cols > 1:
                cap_cx = (ox0 + ox1) / 2
                cap_w = ox1 - ox0
                if cap_w > pw * 0.55:           # spanning caption → full page
                    shade_x0, shade_x1 = 0.0, pw
                elif cap_cx < mid_x:            # left column
                    shade_x0, shade_x1 = 0.0, mid_x
                else:                           # right column
                    shade_x0, shade_x1 = mid_x, pw
            else:
                shade_x0, shade_x1 = ox0, ox1
            prev_caps_0.append((shade_x0, shade_x1, oy1 + 2))

    def _draw_shadowed(rx0: float, rx1: float, ry0: float) -> bool:
        """Return True when a drawing is attributable to a neighbouring figure.

        Uses a minimal 1 pt x-overlap threshold: any drawing that overlaps at
        all with a previous figure's caption x-range (and is above that
        caption's bottom edge) is treated as belonging to that figure.  This is
        intentionally more aggressive than the spanning-detection shadow check
        in _column_x_bounds (which uses 15%) because here we want to protect
        every drawing that could be from an adjacent figure's content area.
        """
        for pcx0, pcx1, pcy_bot in prev_caps_0:
            if ry0 >= pcy_bot:
                continue
            x_ov = min(rx1, pcx1) - max(rx0, pcx0)
            if x_ov > 1:   # any meaningful overlap → belongs to neighbour
                return True
        return False

    draw_top = cy0
    if page is not None:
        try:
            for drw in page.get_drawings():
                rect = drw.get("rect")
                if rect is None:
                    continue
                # Must be above caption start, within column, and within search range.
                # Skip drawings whose bottom edge is at or below the caption start
                # UNLESS the drawing is a large bounding box that spans the
                # figure+caption boundary (e.g. outer figure frame) — in that
                # case only its top edge matters.
                if cy0 - rect.y0 > max_dist:
                    continue
                if rect.y1 >= cy0:
                    # Allow large spanning boxes (outer figure frame) even if they
                    # end past the caption; require area > 500 pt² and top clearly
                    # above caption.  Narrow connectors / arrows are excluded.
                    if rect.y0 >= cy0 - 20 or rect.width * rect.height < 500:
                        continue
                    continue
                if rect.x1 <= col_x0 or rect.x0 >= col_x1:
                    continue
                # Skip effectively zero-dimension path segments (PDF artefacts).
                # A very small threshold (0.5 pt) on the shorter dimension skips
                # zero-width vertical connectors and zero-height horizontal points
                # while keeping thin rules, axis ticks, and bounding boxes that
                # are legitimate parts of figure content.
                if min(rect.width, rect.height) < 0.5:
                    continue
                # Skip if this drawing is more likely from a different figure.
                if _draw_shadowed(rect.x0, rect.x1, rect.y0):
                    continue
                if rect.y0 < draw_top:
                    draw_top = rect.y0
        except Exception:
            pass   # older PyMuPDF versions may not support get_drawings

    if draw_top < cy0 - 30:     # found substantial drawings above the caption
        # Extend draw_top upward to include any raster (embedded image) blocks
        # that sit above the topmost vector drawing in the same column.  Hybrid
        # figures (raster content + vector annotations) would otherwise be
        # cropped at the vector annotation boundary, missing the raster portion.
        for ob in blocks:
            if ob["type"] != 1:
                continue
            ox0, oy0, ox1, oy1 = ob["bbox"]
            if oy1 > cy0 + 5 or cy0 - oy0 > max_dist:
                continue
            img_cx = (ox0 + ox1) / 2
            if not (col_x0 - 15 <= img_cx <= col_x1 + 15):
                continue
            if _draw_shadowed(ox0, ox1, oy0):
                continue
            if oy0 < draw_top:
                draw_top = oy0

        # Before applying the standard 5 pt buffer, check whether there is a
        # text block (header, caption, section title, …) ending just above
        # draw_top.  If so, clip to just below that text to avoid capturing
        # page headers or the bottom of an adjacent figure's caption.
        text_above_bottom: float = 0.0
        for ob in blocks:
            if ob["type"] != 0:
                continue
            _ox0, _oy0, _ox1, oy1 = ob["bbox"]
            if oy1 > draw_top or draw_top - oy1 > 20:
                continue
            txt = "".join(
                s["text"] for ln in ob.get("lines", []) for s in ln.get("spans", [])
            )
            if len(txt.split()) >= 2:
                text_above_bottom = max(text_above_bottom, oy1)
        if text_above_bottom > 0:
            return float(max(0.0, max(text_above_bottom + 1, draw_top - 2)))
        return max(0.0, draw_top - 5)

    # --- Pass 1: embedded raster images ---
    img_candidates: list[tuple[float, float, float, float]] = []
    for ob in blocks:
        if ob["type"] != 1:
            continue
        ox0, oy0, ox1, oy1 = ob["bbox"]
        if oy1 > cy0 + 5 or cy0 - oy0 > max_dist:
            continue
        img_cx = (ox0 + ox1) / 2
        if not (col_x0 - 15 <= img_cx <= col_x1 + 15):
            continue
        img_candidates.append((ox0, oy0, ox1, oy1))

    if img_candidates:
        topmost = min(oy0 for _, oy0, _, _ in img_candidates)
        total_area = sum((x1 - x0) * (y1 - y0) for x0, y0, x1, y1 in img_candidates)
        search_area = col_width * (cy0 - topmost) + 1
        if total_area / search_area >= _MIN_IMG_FILL:
            return topmost

    # --- Pass 2: nearest body-text boundary above the caption ---
    # Body text: block is "wide" (> 25 % of column) and has >= 3 words.
    # This reliably distinguishes paragraphs from axis labels / legend text.
    body_bottom = 0.0          # bottom edge of the nearest qualifying block
    prev_cap_bottom = 0.0      # bottom edge of the nearest figure caption above

    for ob in blocks:
        if ob["type"] != 0:
            continue
        ox0, oy0, ox1, oy1 = ob["bbox"]
        if oy1 >= cy0:
            continue
        if cy0 - oy1 > max_dist:
            continue
        # Horizontal filter: must be roughly in this column.
        ob_cx = (ox0 + ox1) / 2
        if not (col_x0 - 30 <= ob_cx <= col_x1 + 30):
            continue

        text = "".join(
            span["text"]
            for line in ob.get("lines", [])
            for span in line.get("spans", [])
        )

        # Other figure captions act as hard upper bounds.
        if _FIG_RE.search(text):
            if oy1 > prev_cap_bottom:
                prev_cap_bottom = oy1 + 2
            continue

        blk_w = ox1 - ox0
        word_count = len(text.split())
        if blk_w > col_width * 0.25 and word_count >= 3:
            if oy1 > body_bottom:
                body_bottom = oy1 + 2

    # Take the higher (lower on page, i.e. larger y) of body_bottom and prev_cap_bottom.
    candidate = max(body_bottom, prev_cap_bottom)

    # Enforce a minimum figure height of 60 pt.  PDF y increases downward, so
    # "above the caption" means smaller y.  If the candidate is too close to
    # cy0 (figure region too thin), extend it upward (decrease y).
    min_fig_height = 60.0
    if candidate > 0 and cy0 - candidate < min_fig_height:
        # Go up to min_fig_height + 200 pt above caption, but no higher than
        # the previous figure caption bottom (larger y = lower on page = safe).
        extended = max(0.0, cy0 - (min_fig_height + 200))
        candidate = max(extended, prev_cap_bottom)   # don't cross prev caption

    if candidate > 0:
        return candidate

    # --- Pass 3: fixed 300 pt fallback ---
    return max(0.0, cy0 - 300)


# Keywords used by _classify_by_caption.  Ordered from most- to least-specific
# within each category to reduce cross-category false positives.
#
# Removed from _ARCH_KW: "system", "network", "design", "structure" — these
# appear in non-architecture contexts too often (e.g. "system latency",
# "neural network layers", "experimental design", "data structure").
_ARCH_KW = (
    "architecture", "overview", "framework", "pipeline",
    "workflow", "flowchart", "diagram", "module", "routing", "topology",
)
# "breakdown" added (latency/memory breakdowns are performance analyses).
_EVAL_KW = (
    "accuracy", "performance", "result", "comparison", "benchmark",
    "evaluation", "throughput", "latency", "speedup", "memory footprint",
    "tradeoff", "trade-off", "breakdown",
)
# PRIORITY_EVAL_KW are unambiguous evaluation markers that should win even
# when an architecture keyword also appears in the same caption (e.g. an
# ablation study that mentions "pipeline" as a method component).
_PRIORITY_EVAL_KW = (
    "ablation", "perplexity", "f1 score", "top-1", "top-5",
)
_CONCEPT_KW = (
    "illustration", "intuition", "concept", "scenario", "case study", "visuali",
)
# Background / related-work figures: analysis of prior or existing systems.
_BACKGROUND_KW = (
    "prior work", "existing method", "existing approach",
    "motivation", "limitation", "bottleneck", "challenge", "background",
    "related work",
)
# Dataset or data-sample figures.
_DATASET_KW = (
    "dataset", "data sample", "annotation", "data distribution",
    "examples from", "statistics of",
)


def _classify_by_caption(caption: str) -> str:
    """Return a rough figure type from caption keywords (6-class taxonomy).

    Priority order:
    1. PRIORITY_EVAL_KW — unambiguous evaluation markers (ablation, etc.) that
       should override architecture keywords appearing in the same caption.
    2. background_related_work
    3. dataset_example
    4. architecture_flowchart
    5. evaluation_plot
    6. conceptual_illustration
    7. other

    VLM classification (classify_figures) can override this initial label when
    finer accuracy is required.
    """
    cap = caption.lower()
    if any(k in cap for k in _PRIORITY_EVAL_KW):
        return "evaluation_plot"
    if any(k in cap for k in _BACKGROUND_KW):
        return "background_related_work"
    if any(k in cap for k in _DATASET_KW):
        return "dataset_example"
    if any(k in cap for k in _ARCH_KW):
        return "architecture_flowchart"
    if any(k in cap for k in _EVAL_KW):
        return "evaluation_plot"
    if any(k in cap for k in _CONCEPT_KW):
        return "conceptual_illustration"
    return "other"


def _has_visual_content(
    fig_top: float,
    cy0: float,
    col_x0: float,
    col_x1: float,
    blocks: list,
    page: Any,
) -> bool:
    """Return True if the figure area (above caption) contains actual graphics.

    Checks for raster image blocks and vector drawing paths.  When neither is
    present the region is likely a mis-detected caption or a text block, and
    should be skipped.
    """
    if cy0 - fig_top < 20:   # too thin to judge
        return True           # let the normal size checks handle it

    # Raster image block inside the figure region?
    for blk in blocks:
        if blk["type"] != 1:
            continue
        bx0, by0, bx1, by1 = blk["bbox"]
        if bx1 > col_x0 and bx0 < col_x1 and by1 > fig_top and by0 < cy0:
            return True

    # Vector drawing paths inside the figure region.
    # Require at least 3 drawings that are both ≥ 5 pt wide AND ≥ 5 pt tall
    # to avoid false positives from thin table rules / page borders (which
    # commonly have zero or near-zero height/width and would otherwise trigger
    # a false "has visual content" result on text-heavy pages).
    substantial = 0
    try:
        for drw in page.get_drawings():
            rect = drw.get("rect")
            if rect is None:
                continue
            if not (rect.x1 > col_x0 and rect.x0 < col_x1
                    and rect.y1 > fig_top and rect.y0 < cy0):
                continue
            if rect.width >= 5 and rect.height >= 5:
                substantial += 1
                if substantial >= 3:
                    return True
    except Exception:
        pass   # older PyMuPDF versions may not support get_drawings

    return False


def _extract_figures(doc: Any, paper_id: str, out_dir: Path, max_figures: int) -> list[dict]:
    """Locate 'Figure N:' captions in doc, render region as PNG, and embed as SVG wrapper."""
    import fitz  # type: ignore[import]

    figures: list[dict] = []
    seen: set[int] = set()
    scale = 3

    for page_num in range(len(doc)):
        if len(figures) >= max_figures:
            break
        page = doc[page_num]
        pw, ph = page.rect.width, page.rect.height
        blocks = page.get_text("dict")["blocks"]

        # Detect single- vs. two-column layout for this page.
        num_cols, mid_x = _detect_columns(blocks, pw)

        for blk in blocks:
            if blk["type"] != 0:
                continue
            text = "".join(
                span["text"]
                for line in blk["lines"]
                for span in line["spans"]
            )
            m = _FIG_RE.search(text)
            if not m:
                continue
            fig_num = int(m.group(1))
            if fig_num in seen:
                continue

            # Fast filter: if the match is far into the text block it is a
            # mid-sentence reference ("As shown in Figure 12, …") rather than
            # an actual caption that starts the block.  Skip without touching
            # 'seen' so the real caption on a later page can still be found.
            if m.start() > 50:
                continue

            seen.add(fig_num)

            caption = text[m.start():].strip()
            cx0, cy0, cx1, cy1 = blk["bbox"]

            # Determine x-bounds (detects spanning figures automatically).
            col_x0, col_x1 = _column_x_bounds(
                cx0, cx1, num_cols, mid_x, pw, blocks=blocks, cy0=cy0, page=page
            )

            # Estimate where the figure starts above the caption.
            fig_top = _figure_top(
                cy0, cy1, col_x0, col_x1, blocks, page=page,
                num_cols=num_cols, mid_x=mid_x, pw=pw,
            )

            margin = 5
            crop = fitz.Rect(
                max(0.0, col_x0 - margin),
                max(0.0, fig_top - margin),
                min(pw, col_x1 + margin),
                min(ph, cy1 + margin),
            )

            # Skip degenerate crops (caption-only slivers).
            if crop.height < 30 or crop.width < 30:
                continue

            # Skip crops whose figure area (above caption) contains no graphics —
            # these are usually long captions or body-text blocks mis-detected as
            # figures (e.g. a figure with an unusually long multi-sentence caption
            # that gets confused with a new figure reference in the text).
            if not _has_visual_content(fig_top, cy0, col_x0, col_x1, blocks, page):
                logger.debug(
                    "Skipping Fig {} p{}: no graphics found above caption",
                    fig_num, page_num + 1,
                )
                continue

            mat = fitz.Matrix(scale, scale)
            pix = page.get_pixmap(matrix=mat, clip=crop)

            # Skip suspiciously small rendered images.
            if pix.width < 60 or pix.height < 60:
                continue

            idx = len(figures) + 1
            base_name = f"fig{idx:02d}_p{page_num + 1}"
            png_path = out_dir / f"{base_name}.png"
            pix.save(str(png_path))

            # Wrap the rasterised PNG as an SVG so downstream tools get a
            # vector-container that can be opened and re-layered in Inkscape.
            w_pt = crop.width
            h_pt = crop.height
            svg_content = (
                f'<svg xmlns="http://www.w3.org/2000/svg" '
                f'xmlns:xlink="http://www.w3.org/1999/xlink" '
                f'width="{w_pt:.1f}" height="{h_pt:.1f}" '
                f'viewBox="0 0 {w_pt:.1f} {h_pt:.1f}">\n'
                f'  <image href="{png_path.name}" '
                f'width="{w_pt:.1f}" height="{h_pt:.1f}"/>\n'
                f'</svg>\n'
            )
            svg_path = out_dir / f"{base_name}.svg"
            svg_path.write_text(svg_content, encoding="utf-8")

            figures.append({
                "paper_id": paper_id,
                "fig_index": idx,
                "fig_num": fig_num,
                "page": page_num + 1,
                "caption": caption,
                # Heuristic label from caption keywords; overridden by VLM
                # classify_figures when finer accuracy is needed.
                "figure_type": _classify_by_caption(caption),
                "png_path": str(png_path),
                "svg_path": str(svg_path),
                "width_px": pix.width,
                "height_px": pix.height,
            })

            if len(figures) >= max_figures:
                break

    return figures


# ---------------------------------------------------------------------------
# CSV helpers
# ---------------------------------------------------------------------------

_PER_PAPER_FIELDS = [
    "fig_index", "fig_num", "page", "caption",
    "figure_type", "png_path", "svg_path", "width_px", "height_px",
]

_CATALOG_FIELDS = [
    "paper_id", "paper_title", "paper_year", "paper_citations",
    "fig_index", "fig_num", "page", "caption",
    "figure_type", "png_path", "svg_path", "in_pptx",
]


def _write_figures_csv(path: Path, figures: list[dict]) -> None:
    """Write per-paper figure metadata to a CSV file."""
    with open(path, "w", newline="", encoding="utf-8") as f:
        w = csv.DictWriter(f, fieldnames=_PER_PAPER_FIELDS, extrasaction="ignore")
        w.writeheader()
        for fig in figures:
            row = {k: fig.get(k, "") for k in _PER_PAPER_FIELDS}
            row["caption"] = (row["caption"] or "")[:300]
            w.writerow(row)


def _write_catalog_csv(
    path: Path, figures: list[dict], selected_pngs: set
) -> None:
    """Write a consolidated figure catalogue across all papers."""
    with open(path, "w", newline="", encoding="utf-8") as f:
        w = csv.DictWriter(f, fieldnames=_CATALOG_FIELDS, extrasaction="ignore")
        w.writeheader()
        for fig in figures:
            row = {k: fig.get(k, "") for k in _CATALOG_FIELDS}
            row["caption"] = (row["caption"] or "")[:300]
            row["in_pptx"] = "yes" if fig.get("png_path") in selected_pngs else "no"
            w.writerow(row)


# ---------------------------------------------------------------------------
# PPTX export helper
# ---------------------------------------------------------------------------

def _build_pptx(
    figures: list[dict],
    out_path: str,
    topic: str,
    figure_type: str = "",
) -> None:
    """Write a styled PPTX: centered cover slide + one figure per content slide."""
    from pptx import Presentation  # type: ignore[import]
    from pptx.util import Inches, Pt  # type: ignore[import]
    from pptx.dml.color import RGBColor  # type: ignore[import]
    from pptx.enum.text import PP_ALIGN  # type: ignore[import]

    _DARK = RGBColor(0x1B, 0x1B, 0x2F)
    _BLUE = RGBColor(0x4A, 0x90, 0xD9)
    _WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    _LBLUE = RGBColor(0x7E, 0xB8, 0xFF)
    _GREY = RGBColor(0x88, 0x88, 0xAA)
    _CAPGREY = RGBColor(0x44, 0x44, 0x44)
    _SUBGREY = RGBColor(0x88, 0x88, 0x88)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]   # completely blank

    # ── Cover slide ────────────────────────────────────────────────────────────
    cover = prs.slides.add_slide(blank)

    bg = cover.background.fill
    bg.solid()
    bg.fore_color.rgb = _DARK

    def _bar(slide: Any, top: float, h: float = 0.07) -> None:
        s = slide.shapes.add_shape(1, Inches(0), Inches(top), Inches(13.33), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = _BLUE
        s.line.fill.background()

    _bar(cover, 0)
    _bar(cover, 7.43)

    def _txb(slide: Any, top: float, h: float, size: int, text: str,
             color: RGBColor, bold: bool = False, space_before: int = 0) -> None:
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(top), Inches(10.33), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        if space_before:
            p.space_before = Pt(space_before)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color

    _txb(cover, 2.1, 0.85, 32, f"Topic: {topic}", _WHITE, bold=True)
    # Figure type line always shows; defaults to "All" when no filter is applied.
    _txb(cover, 3.2, 0.65, 22,
         f"Figure type: {figure_type or 'All'}", _LBLUE)
    _txb(cover, 5.5, 0.5, 13,
         f"{len(figures)} reference figure(s)  ·  ClawPhD 🐈", _GREY)

    # ── Content slides ─────────────────────────────────────────────────────────
    # Layout (7.5" tall slide):
    #   0.25" – 5.85"  image area  (5.6" max height)
    #   6.0"  – 6.85"  caption     (0.85", Pt 13)
    #   6.9"  – 7.45"  paper info  (0.55", Pt 11)
    _IMG_TOP  = Inches(0.25)
    _IMG_H    = Inches(5.6)
    _CAP_TOP  = Inches(6.0)
    _CAP_H    = Inches(0.85)
    _INFO_TOP = Inches(6.9)
    _INFO_H   = Inches(0.55)
    _SIDE_M   = Inches(0.4)
    _TXB_W    = prs.slide_width - _SIDE_M * 2

    for fig in figures:
        png = fig.get("png_path", "")
        if not png or not Path(png).exists():
            continue
        slide = prs.slides.add_slide(blank)

        # Insert image, then enforce both height AND width limits.
        try:
            pic = slide.shapes.add_picture(png, _SIDE_M, _IMG_TOP, height=_IMG_H)
            max_w = prs.slide_width - _SIDE_M * 2
            if pic.width > max_w:
                # Scale down proportionally so it fits within the slide.
                scale = max_w / pic.width
                pic.width = int(max_w)
                pic.height = int(pic.height * scale)
            # Centre horizontally.
            pic.left = int((prs.slide_width - pic.width) / 2)
        except Exception:
            pass

        # Caption line
        caption = (fig.get("caption") or "")[:200]
        fig_type = fig.get("figure_type", "")
        label = f"[{fig_type}]  " if fig_type and fig_type not in ("", "other") else ""
        cap_tb = slide.shapes.add_textbox(_SIDE_M, _CAP_TOP, _TXB_W, _CAP_H)
        cap_tf = cap_tb.text_frame
        cap_tf.word_wrap = True
        cp = cap_tf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = label + caption
        cr.font.size = Pt(13)
        cr.font.color.rgb = _CAPGREY

        # Paper info line
        paper_info = f"{fig.get('paper_title', '')} ({fig.get('paper_year', '')})"
        if paper_info.strip("() "):
            pi_tb = slide.shapes.add_textbox(_SIDE_M, _INFO_TOP, _TXB_W, _INFO_H)
            pi_tf = pi_tb.text_frame
            pip = pi_tf.paragraphs[0]
            pip.alignment = PP_ALIGN.CENTER
            pir = pip.add_run()
            pir.text = paper_info
            pir.font.size = Pt(11)
            pir.font.color.rgb = _SUBGREY

    prs.save(out_path)


# ---------------------------------------------------------------------------
# Tool 1 – search_influential_papers
# ---------------------------------------------------------------------------

class SearchInfluentialPapersTool(Tool):
    """Search for influential papers on a topic via Semantic Scholar (+ OpenAlex fallback)."""

    name = "search_influential_papers"
    description = (
        "Search for influential academic papers on a given topic. "
        "Uses Semantic Scholar REST API (no key required). "
        "Returns a list of papers with titles, citation counts, years, and open-access PDF URLs."
    )
    parameters = {
        "type": "object",
        "properties": {
            "topic": {
                "type": "string",
                "description": "English-language topic or keywords to search for.",
            },
            "num_papers": {
                "type": "integer",
                "description": "Maximum number of papers to return (default: 5).",
            },
            "min_citations": {
                "type": "integer",
                "description": "Minimum citation count filter (default: 50).",
            },
            "year_start": {
                "type": "integer",
                "description": "Only include papers published from this year onward (default: 2017).",
            },
            "expand_citations": {
                "type": "boolean",
                "description": (
                    "If true, also fetch highly-cited references of the seed results "
                    "to widen the pool (default: false, saves API quota)."
                ),
            },
        },
        "required": ["topic"],
    }

    def __init__(self, api_key: str | None = None):
        self._api_key = api_key or None

    async def execute(
        self,
        topic: str,
        num_papers: int = 5,
        min_citations: int = 50,
        year_start: int = 2017,
        expand_citations: bool = False,
        **kwargs: Any,
    ) -> str:
        params = {
            "query": topic,
            "fields": _S2_FIELDS,
            "limit": min(num_papers * 4, 100),
        }

        s2_error: str | None = None
        pool: list[dict] = []

        # --- Semantic Scholar seed search ---
        try:
            async with _s2_client(self._api_key) as client:
                seed_resp = await _s2_get(client, _S2_SEARCH, params=params)
                pool = list(seed_resp.json().get("data") or [])
            logger.info("S2 seed search '{}' → {} papers", topic, len(pool))

            # Optional citation expansion (costs extra API calls)
            if expand_citations and pool:
                extra: list[dict] = []
                for paper in pool[:3]:
                    pid = paper.get("paperId")
                    if not pid:
                        continue
                    await asyncio.sleep(1)
                    try:
                        async with _s2_client(self._api_key) as client:
                            ref_resp = await _s2_get(
                                client,
                                _S2_REFS.format(pid=pid),
                                params={"fields": _S2_FIELDS, "limit": 20},
                            )
                        for item in ref_resp.json().get("data") or []:
                            cited = item.get("citedPaper") or {}
                            if cited:
                                extra.append(cited)
                    except Exception as e:
                        logger.warning("S2 ref expansion failed for {}: {}", pid, e)
                pool = pool + extra
        except Exception as e:
            s2_error = str(e)
            logger.warning("S2 search failed ({}), trying OpenAlex fallback", s2_error)

        # --- Filter and deduplicate S2 results ---
        results: list[dict] = []
        seen_ids: set[str] = set()

        for paper in pool:
            pid = paper.get("paperId") or ""
            if pid in seen_ids:
                continue
            seen_ids.add(pid)

            cit = paper.get("citationCount") or 0
            year = paper.get("year") or 0
            if cit < min_citations or year < year_start:
                continue

            ext = paper.get("externalIds") or {}
            arxiv_id = ext.get("ArXiv")
            pdf_url: str | None = None
            oa = paper.get("openAccessPdf") or {}
            pdf_url = oa.get("url")
            if not pdf_url and arxiv_id:
                pdf_url = f"https://arxiv.org/pdf/{arxiv_id}"

            if not pdf_url:
                continue

            results.append({
                "paper_id": arxiv_id or _slugify(paper.get("title", pid))[:20],
                "arxiv_id": arxiv_id,
                "title": paper.get("title") or "",
                "year": year,
                "citation_count": cit,
                "pdf_url": pdf_url,
                "source": "semantic_scholar",
            })
            if len(results) >= num_papers:
                break

        # --- OpenAlex fallback when S2 yields nothing ---
        if not results and (s2_error or not pool):
            logger.info("Falling back to OpenAlex for '{}'", topic)
            try:
                results = await _openalex_search(topic, num_papers, min_citations, year_start)
                if results:
                    logger.info("OpenAlex returned {} papers for '{}'", len(results), topic)
            except Exception as e:
                logger.warning("OpenAlex fallback also failed: {}", e)
                if s2_error:
                    return (
                        f"Error: both Semantic Scholar ({s2_error}) and OpenAlex ({e}) "
                        "failed. Please try again in a few minutes."
                    )

        if not results:
            return (
                f"Error: no open-access papers found for '{topic}' "
                f"(min_citations={min_citations}, year_start={year_start}). "
                "Try lowering min_citations or year_start."
            )

        # Sort by citation count descending
        results.sort(key=lambda p: p.get("citation_count", 0), reverse=True)
        results = results[:num_papers]
        logger.info("Returning {} papers for '{}'", len(results), topic)
        return json.dumps(results, indent=2, ensure_ascii=False)


# ---------------------------------------------------------------------------
# Tool 2 – extract_paper_figures
# ---------------------------------------------------------------------------

class ExtractPaperFiguresTool(Tool):
    """Download a paper PDF and extract all labelled figures as PNG + SVG files."""

    name = "extract_paper_figures"
    description = (
        "Download an academic paper PDF and extract every labelled figure "
        "(Figure N / Fig. N captions) as a PNG raster and an SVG wrapper. "
        "Returns a list of figure metadata including file paths."
    )
    parameters = {
        "type": "object",
        "properties": {
            "paper_id": {
                "type": "string",
                "description": "Unique identifier for the paper (arxiv_id preferred, else slug).",
            },
            "pdf_url": {
                "type": "string",
                "description": "Direct URL to the paper PDF.",
            },
            "paper_title": {
                "type": "string",
                "description": "Paper title (stored in figure metadata).",
            },
            "paper_year": {
                "type": "integer",
                "description": "Publication year.",
            },
            "paper_citations": {
                "type": "integer",
                "description": "Citation count.",
            },
            "max_figures": {
                "type": "integer",
                "description": "Maximum figures to extract per paper (default: 20).",
            },
            "output_dir": {
                "type": "string",
                "description": "Parent directory for extracted figures (default: <workspace>/outputs/figure_refs).",
            },
        },
        "required": ["paper_id", "pdf_url"],
    }

    def __init__(self, workspace: Path, allowed_dir: Path | None = None):
        self._workspace = workspace
        self._allowed_dir = allowed_dir

    async def execute(
        self,
        paper_id: str,
        pdf_url: str,
        paper_title: str = "",
        paper_year: int = 0,
        paper_citations: int = 0,
        max_figures: int = 20,
        output_dir: str | None = None,
        **kwargs: Any,
    ) -> str:
        try:
            import fitz  # type: ignore[import]  # noqa: F401
        except ImportError:
            return "Error: PyMuPDF is required. Run: pip install PyMuPDF"

        try:
            import httpx
        except ImportError:
            return "Error: httpx is required. Run: pip install httpx"

        out_base = Path(output_dir) if output_dir else self._workspace / "outputs" / "figure_refs"
        paper_dir = out_base / paper_id
        paper_dir.mkdir(parents=True, exist_ok=True)

        # Remove stale figure files from previous extraction runs so that the
        # results directory only contains what the current code produces.
        # paper.pdf and figures.csv are regenerated later and do not need to
        # be deleted here; only the rendered PNG/SVG crops are cleared.
        for old_fig in list(paper_dir.glob("fig*.png")) + list(paper_dir.glob("fig*.svg")):
            try:
                old_fig.unlink()
            except Exception:
                pass

        # --- Download PDF ---
        pdf_path = paper_dir / "paper.pdf"
        if not pdf_path.exists():
            logger.info("Downloading PDF for {}: {}", paper_id, pdf_url)
            try:
                async with httpx.AsyncClient(timeout=60.0, follow_redirects=True) as client:
                    resp = await client.get(pdf_url)
                    resp.raise_for_status()
                    content_type = resp.headers.get("content-type", "")
                    if "pdf" not in content_type and not resp.content[:4] == b"%PDF":
                        return f"Error: URL did not return a PDF (content-type: {content_type})"
                    pdf_path.write_bytes(resp.content)
            except Exception as e:
                return f"Error: failed to download PDF from {pdf_url}: {e}"
        else:
            logger.info("PDF already cached for {}", paper_id)

        # --- Extract figures ---
        try:
            import fitz  # type: ignore[import]

            doc = fitz.open(str(pdf_path))
            figures = _extract_figures(doc, paper_id, paper_dir, max_figures)
            doc.close()
        except Exception as e:
            logger.exception("Figure extraction failed for {}", paper_id)
            return f"Error: figure extraction failed for {paper_id}: {e}"

        # Attach paper metadata to every figure record
        for fig in figures:
            fig["paper_title"] = paper_title
            fig["paper_year"] = paper_year
            fig["paper_citations"] = paper_citations

        # Write per-paper CSV catalogue (useful for debugging extraction quality)
        csv_path = paper_dir / "figures.csv"
        _write_figures_csv(csv_path, figures)

        logger.info("Extracted {} figures from {} → {}", len(figures), paper_id, paper_dir)
        return json.dumps(
            {
                "paper_id": paper_id,
                "paper_title": paper_title,
                "figures_extracted": len(figures),
                "output_dir": str(paper_dir),
                "catalog_csv": str(csv_path),
                "figures": figures,
            },
            indent=2,
            ensure_ascii=False,
        )


# ---------------------------------------------------------------------------
# Tool 3 – classify_figures
# ---------------------------------------------------------------------------

class ClassifyFiguresTool(Tool):
    """Classify extracted figures by type using the configured LLM vision provider."""

    name = "classify_figures"
    description = (
        "Classify a list of extracted paper figures by type "
        "(architecture_flowchart / evaluation_plot / conceptual_illustration / table / other) "
        "using the already-configured LLM provider — no extra API key required. "
        "Adds a 'figure_type' field to each figure dict and returns the updated list."
    )
    parameters = {
        "type": "object",
        "properties": {
            "figures": {
                "type": "array",
                "description": "List of figure dicts from extract_paper_figures.",
                "items": {"type": "object"},
            },
        },
        "required": ["figures"],
    }

    _TYPES = [
        "architecture_flowchart",
        "evaluation_plot",
        "conceptual_illustration",
        "table",
        "other",
    ]
    _PROMPT = (
        "You are an academic figure classifier.\n"
        "Classify the image into exactly one of these categories:\n"
        "  architecture_flowchart – system diagram, pipeline, framework overview\n"
        "  evaluation_plot        – bar chart, line plot, ablation table rendered as image\n"
        "  conceptual_illustration – motivation figure, teaser, concept sketch\n"
        "  table                  – data table (even if rendered as image)\n"
        "  other                  – anything else\n\n"
        'Return strict JSON only: {"figure_type": "<category>"}'
    )

    def __init__(self, vlm_provider: Any = None):
        self._vlm = vlm_provider

    async def execute(self, figures: list[dict], **kwargs: Any) -> str:
        if not self._vlm:
            return "Error: VLM provider not configured."
        if not figures:
            return json.dumps([], ensure_ascii=False)

        try:
            from PIL import Image  # type: ignore[import]
        except ImportError:
            return "Error: Pillow is required. Run: pip install Pillow"

        classified: list[dict] = []
        for fig in figures:
            png = fig.get("png_path", "")
            fig_type = "other"
            if png and Path(png).exists():
                try:
                    img = Image.open(png).convert("RGB")
                    max_dim = 1024
                    if max(img.size) > max_dim:
                        ratio = max_dim / max(img.size)
                        resample = getattr(Image, "LANCZOS", getattr(Image, "ANTIALIAS", None))
                        img = img.resize(
                            (int(img.width * ratio), int(img.height * ratio)),
                            resample,
                        )
                    resp = await self._vlm.generate(
                        prompt=self._PROMPT,
                        images=[img],
                        temperature=0.0,
                        response_format="json",
                    )
                    raw = resp.strip()
                    # Strip markdown fences if present
                    if raw.startswith("```"):
                        raw = re.sub(r"^```[a-z]*\n?", "", raw, flags=re.IGNORECASE)
                        raw = raw.rstrip("` \n")
                    parsed = json.loads(raw)
                    candidate = parsed.get("figure_type", "other")
                    if candidate in self._TYPES:
                        fig_type = candidate
                except Exception as e:
                    logger.warning("Classification failed for {}: {}", png, e)
            classified.append({**fig, "figure_type": fig_type})

        return json.dumps(classified, indent=2, ensure_ascii=False)


# ---------------------------------------------------------------------------
# Tool 4 – export_figure_reference
# ---------------------------------------------------------------------------

class ExportFigureReferenceTool(Tool):
    """Filter figures by type and export them to PPTX and/or return SVG paths."""

    name = "export_figure_reference"
    description = (
        "Filter a list of classified figures by type and export as an editable "
        "PowerPoint (PPTX) and/or return SVG file paths. "
        "Each slide in the PPTX contains one figure with its caption."
    )
    parameters = {
        "type": "object",
        "properties": {
            "figures": {
                "type": "array",
                "description": "List of figure dicts (from extract_paper_figures or classify_figures).",
                "items": {"type": "object"},
            },
            "figure_type_filter": {
                "type": "array",
                "description": (
                    "Only include figures whose 'figure_type' matches one of these values. "
                    "Omit or pass [] to include all figures. "
                    "Values: architecture_flowchart | evaluation_plot | conceptual_illustration | table | other"
                ),
                "items": {"type": "string"},
            },
            "output_format": {
                "type": "string",
                "description": "What to produce: 'pptx', 'svg', or 'both' (default: 'both').",
                "enum": ["pptx", "svg", "both"],
            },
            "slide_title": {
                "type": "string",
                "description": "Title shown on the PPTX cover slide (default: 'Figure References').",
            },
            "output_path": {
                "type": "string",
                "description": (
                    "Full path for the PPTX file. "
                    "If omitted, saved to <workspace>/outputs/figure_refs/reference_pack_<title>_<ts>.pptx."
                ),
            },
        },
        "required": ["figures"],
    }

    def __init__(self, workspace: Path):
        self._workspace = workspace

    async def execute(
        self,
        figures: list[dict],
        figure_type_filter: list[str] | None = None,
        output_format: str = "both",
        slide_title: str = "Figure References",
        output_path: str | None = None,
        **kwargs: Any,
    ) -> str:
        if not figures:
            return "Error: no figures provided."

        # Pre-generate shared slug + timestamp so PPTX and CSV always match.
        out_dir = self._workspace / "outputs" / "figure_refs"
        out_dir.mkdir(parents=True, exist_ok=True)
        safe_title = _slugify(slide_title)[:40] or "refs"
        ts = _timestamp()
        base_name = f"reference_pack_{safe_title}_{ts}"

        # --- Filter ---
        active_filter = [f for f in (figure_type_filter or []) if f]
        figure_type_label = ", ".join(active_filter) if active_filter else ""
        if active_filter:
            selected = [f for f in figures if f.get("figure_type") in active_filter]
            filter_note: str | None = (
                f"Filtered {len(figures)} → {len(selected)} figures "
                f"matching types: {active_filter}"
            )
            # Fall back to all figures if filter yields nothing
            if not selected:
                selected = figures
                filter_note = (
                    f"No figures matched filter {active_filter}; "
                    "returning all figures instead."
                )
        else:
            selected = figures
            filter_note = None

        result: dict[str, Any] = {
            "total_extracted": len(figures),
            "selected_count": len(selected),
        }
        if filter_note:
            result["note"] = filter_note

        # --- SVG paths ---
        if output_format in ("svg", "both"):
            svg_files = [
                f["svg_path"]
                for f in selected
                if f.get("svg_path") and Path(f["svg_path"]).exists()
            ]
            result["svg_files"] = svg_files
            result["svg_count"] = len(svg_files)

        # --- PPTX ---
        pptx_path: str | None = None
        if output_format in ("pptx", "both"):
            pptx_path = output_path or str(out_dir / f"{base_name}.pptx")
            try:
                _build_pptx(selected, pptx_path, topic=slide_title,
                            figure_type=figure_type_label)
                result["pptx_path"] = pptx_path
                result["pptx_slides"] = len(selected)
                logger.info("PPTX saved: {} ({} slides)", pptx_path, len(selected))
            except ImportError:
                result["pptx_error"] = (
                    "python-pptx not installed. Run: pip install python-pptx"
                )
            except Exception as e:
                logger.exception("PPTX build failed")
                result["pptx_error"] = str(e)

        # --- Catalogue CSV (same base-name as PPTX, or standalone with timestamp) ---
        try:
            if pptx_path:
                # Always co-locate with the PPTX — replace extension only.
                catalog_path = Path(pptx_path).with_suffix(".csv")
            else:
                catalog_path = out_dir / f"{base_name}.csv"
            selected_pngs = {f.get("png_path") for f in selected}
            _write_catalog_csv(catalog_path, figures, selected_pngs)
            result["catalog_csv"] = str(catalog_path)
            logger.info("Catalog CSV saved: {}", catalog_path)
        except Exception as e:
            logger.warning("CSV catalogue write failed: {}", e)

        # Note: per-paper figures.csv is written (and remains authoritative) by
        # ExtractPaperFiguresTool.  We intentionally do NOT overwrite it here,
        # because export receives only a subset of figures and would lose the
        # complete per-paper record.  If classify_figures updated figure_type
        # values, those are reflected in the catalog CSV below instead.

        return json.dumps(result, indent=2, ensure_ascii=False)
